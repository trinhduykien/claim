# -*- coding: utf-8 -*-
"""Táº¡o PowerPoint thuyáº¿t trÃ¬nh AI Claim Chatbot V.0.10"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# â”€â”€ Colors â”€â”€
BRAND_BLUE = RGBColor(0x00, 0x2B, 0x70)
BRAND_ORANGE = RGBColor(0xF5, 0x82, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xCC, 0x33, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=WHITE):
    """Fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=18, color=DARK_TEXT, bullet=True, spacing=Pt(8), font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "â€¢ " if bullet else ""
        p.text = prefix + text
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox

def slide_header(slide, title, subtitle=None):
    """Add header band + title."""
    add_rect(slide, 0, 0, W, Inches(1.1), BRAND_BLUE)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6), title,
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.4), subtitle,
                 font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF))
    # Orange accent line
    add_rect(slide, 0, Inches(1.1), W, Pt(3), BRAND_ORANGE)

def slide_footer(slide, page_num):
    """Add footer."""
    add_text(slide, Inches(0.6), Inches(7.0), Inches(6), Inches(0.4),
             "AI Claim Chatbot V.0.10", font_size=10, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
             str(page_num), font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BRAND_BLUE)
add_rect(slide, 0, Inches(2.5), W, Inches(2.5), RGBColor(0x00, 0x1F, 0x50))
add_text(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(1.2),
         "AI Claim Chatbot", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(0.6),
         "Trá»£ lÃ½ áº£o Ä‘Ã¡nh giÃ¡ tiáº¿p nháº­n bá»“i thÆ°á»ng báº£o hiá»ƒm", font_size=22, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
         "PhiÃªn báº£n V.0.10  |  02/07/2026", font_size=16, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)
# Logo text
add_text(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.6),
         "", font_size=28, bold=True, color=BRAND_ORANGE, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(1.2), Inches(2.3), Pt(3), BRAND_ORANGE)

# ============================================================
# SLIDE 2: Váº¥n Ä‘á»
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "Váº¥n Ä‘á»", "Táº¡i sao cáº§n trá»£ lÃ½ AI Ä‘Ã¡nh giÃ¡ bá»“i thÆ°á»ng?")
slide_footer(slide, 2)

add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
    [
        ("ÄÃ¡nh giÃ¡ bá»“i thÆ°á»ng thá»§ cÃ´ng tá»‘n thá»i gian", True),
        ("NhÃ¢n viÃªn pháº£i Ä‘á»c láº¡i há»£p Ä‘á»“ng, Ä‘á»‘i chiáº¿u hÃ³a Ä‘Æ¡n tá»«ng dÃ²ng", False),
        ("Dá»… bá» sÃ³t Ä‘iá»u khoáº£n loáº¡i trá»« giÃ¡n tiáº¿p", False),
        ("Phá»¥ lá»¥c, Ä‘Ã­nh chá»‰nh náº±m á»Ÿ trang khÃ¡c nhau", False),
        ("", False),
        ("KhÃ¡ch hÃ ng chá» Ä‘á»£i lÃ¢u", True),
        ("KhÃ´ng cÃ³ pháº£n há»“i ngay láº­p tá»©c", False),
        ("Há»“ sÆ¡ lÆ°u trá»¯ thá»§ cÃ´ng, khÃ³ tra cá»©u", False),
        ("", False),
        ("Kiáº¿n thá»©c phÃ¢n tÃ¡n", True),
        ("Má»—i Ä‘Ã¡nh giÃ¡ viÃªn hiá»ƒu khÃ¡c nhau vá» Ä‘iá»u khoáº£n", False),
        ("KhÃ´ng cÃ³ tiÃªu chuáº©n chung cho viá»‡c Ä‘á»‘i chiáº¿u", False),
    ], font_size=16, bullet=True)

# Right side - solution preview
add_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(4.8), LIGHT_BG)
add_text(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
         "Giáº£i phÃ¡p: AI Claim Chatbot", font_size=20, bold=True, color=BRAND_BLUE)
add_multi_text(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(4),
    [
        ("Chatbot trá»±c tuyáº¿n 24/7", True),
        ("KhÃ¡ch hÃ ng chat â†’ AI tá»± Ä‘Ã¡nh giÃ¡", False),
        ("", False),
        ("Äá»c há»£p Ä‘á»“ng + hÃ³a Ä‘Æ¡n tá»± Ä‘á»™ng", True),
        ("AI Ä‘á»c áº£nh PDF, trÃ­ch xuáº¥t text, Ä‘á»‘i chiáº¿u", False),
        ("", False),
        ("Xuáº¥t báº£ng kháº¥u trá»« chÃ­nh xÃ¡c", True),
        ("Má»—i má»¥c Ä‘á»™c láº­p, cÃ³ trÃ­ch dáº«n Ä‘iá»u khoáº£n", False),
        ("", False),
        ("Há»“ sÆ¡ lÆ°u tá»± Ä‘á»™ng", True),
        ("JSON + Excel, dá»… tra cá»©u", False),
    ], font_size=16, bullet=True, color=DARK_TEXT)

# ============================================================
# SLIDE 3: Tá»•ng quan há»‡ thá»‘ng
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "Tá»•ng quan há»‡ thá»‘ng", "Kiáº¿n trÃºc 3 táº§ng Map-Reduce-Merge")
slide_footer(slide, 3)

# 3 tiers as boxes
tier_y = Inches(1.8)
tier_h = Inches(1.5)
tier_w = Inches(3.5)
gap = Inches(0.5)

# Tier 1
add_rect(slide, Inches(0.8), tier_y, tier_w, tier_h, BRAND_BLUE)
add_text(slide, Inches(1), Inches(1.95), Inches(3.1), Inches(0.5),
         "TIER 1 â€” MAP", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.4), Inches(3.1), Inches(0.8),
         "Kimi K2.6\nÄá»c áº£nh hÃ³a Ä‘Æ¡n + há»£p Ä‘á»“ng\nTrÃ­ch xuáº¥t text nguyÃªn vÄƒn", font_size=13, color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(4.4), Inches(2.2), Inches(0.4), Inches(0.5), "â†’", font_size=28, bold=True, color=BRAND_ORANGE, align=PP_ALIGN.CENTER)

# Tier 2
add_rect(slide, Inches(4.9), tier_y, tier_w, tier_h, BRAND_ORANGE)
add_text(slide, Inches(5.1), Inches(1.95), Inches(3.1), Inches(0.5),
         "TIER 2 â€” REDUCE", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(5.1), Inches(2.4), Inches(3.1), Inches(0.8),
         "GLM-5.2\nPhÃ¢n tÃ­ch tá»«ng chunk song song\nTrÃ­ch xuáº¥t A/B/C tá»« há»£p Ä‘á»“ng", font_size=13, color=RGBColor(0xFF, 0xF0, 0xDD), align=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(8.5), Inches(2.2), Inches(0.4), Inches(0.5), "â†’", font_size=28, bold=True, color=BRAND_ORANGE, align=PP_ALIGN.CENTER)

# Tier 3
add_rect(slide, Inches(9), tier_y, tier_w, tier_h, BRAND_BLUE)
add_text(slide, Inches(9.2), Inches(1.95), Inches(3.1), Inches(0.5),
         "TIER 3 â€” MERGE", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.2), Inches(2.4), Inches(3.1), Inches(0.8),
         "GLM-5.2\nTá»•ng há»£p + suy luáº­n kháº¥u trá»«\nXuáº¥t báº£ng káº¿t quáº£ cuá»‘i cÃ¹ng", font_size=13, color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

# Below: tech stack
add_text(slide, Inches(0.8), Inches(3.8), Inches(12), Inches(0.5),
         "Tech Stack", font_size=20, bold=True, color=BRAND_BLUE)
add_multi_text(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(3),
    [
        ("Frontend: Streamlit (Python)", False),
        ("AI Vision: Kimi K2.6 (Ollama Cloud)", False),
        ("AI Analysis: GLM-5.2 (Ollama Cloud)", False),
        ("PDF Processing: PyMuPDF (fitz)", False),
        ("Bot Pattern: python-aiml", False),
    ], font_size=15, bullet=True)

add_multi_text(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(3),
    [
        ("Backend: Python 3.14", False),
        ("Threading: ThreadPoolExecutor (song song)", False),
        ("Storage: JSON (claim_logs/) + Markdown", False),
        ("API: Ollama Cloud (REST)", False),
        ("Version: Git + GitHub (tag V.0.10)", False),
    ], font_size=15, bullet=True)

# ============================================================
# SLIDE 4: Luá»“ng hoáº¡t Ä‘á»™ng
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "Luá»“ng hoáº¡t Ä‘á»™ng", "Tá»« khÃ¡ch hÃ ng chat â†’ káº¿t luáº­n bá»“i thÆ°á»ng")
slide_footer(slide, 4)

steps = [
    ("1", "KhÃ¡ch hÃ ng chat", "Má»Ÿ web â†’ bot chÃ o há»i tá»± nhiÃªn (AIML)"),
    ("2", "MÃ´ táº£ sá»± cá»‘", "KhÃ¡ch nÃ³i 'nhÃ  tÃ´i bá»‹ chÃ¡y, cÃ³ mua Combo 360'"),
    ("3", "AI nháº­n diá»‡n sáº£n pháº©m", "Match vá»›i 11 sáº£n pháº©m "),
    ("4", "AI há»i Ä‘Ã¡nh giÃ¡", "Má»—i láº§n 1-2 cÃ¢u, linh hoáº¡t nhÆ° Ä‘Ã¡nh giÃ¡ viÃªn"),
    ("5", "KhÃ¡ch upload hÃ³a Ä‘Æ¡n + há»£p Ä‘á»“ng", "áº¢nh hÃ³a Ä‘Æ¡n + PDF há»£p Ä‘á»“ng"),
    ("6", "AI phÃ¢n tÃ­ch kháº¥u trá»«", "3 táº§ng: Ä‘á»c áº£nh â†’ phÃ¢n tÃ­ch â†’ xuáº¥t báº£ng"),
    ("7", "Káº¿t luáº­n", "âœ… Äá»§ Ä‘iá»u kiá»‡n / âŒ KhÃ´ng Ä‘á»§ Ä‘iá»u kiá»‡n + lÃ½ do"),
    ("8", "LÆ°u há»“ sÆ¡", "JSON + Markdown tá»± Ä‘á»™ng"),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.5 + i * 0.65)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BRAND_ORANGE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.word_wrap = False
    # Title + desc
    add_text(slide, Inches(1.5), y, Inches(3.5), Inches(0.5), title,
             font_size=16, bold=True, color=BRAND_BLUE)
    add_text(slide, Inches(5.2), y, Inches(7), Inches(0.5), desc,
             font_size=14, color=DARK_TEXT)

# ============================================================
# SLIDE 5: Pipeline 3 táº§ng chi tiáº¿t
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "Pipeline AI â€” 3 táº§ng chi tiáº¿t", "Map â†’ Reduce â†’ Merge")
slide_footer(slide, 5)

# Tier 1 detail
add_rect(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(1.6), LIGHT_BG)
add_text(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.4),
         "TIER 1 â€” MAP (Kimi K2.6)", font_size=16, bold=True, color=BRAND_BLUE)
add_multi_text(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(1.2),
    [
        ("Äá»c áº£nh hÃ³a Ä‘Æ¡n: trÃ­ch xuáº¥t tá»«ng dÃ²ng â€” tÃªn má»¥c, SL, Ä‘Æ¡n giÃ¡, thÃ nh tiá»n", False),
        ("Äá»c áº£nh há»£p Ä‘á»“ng: chia thÃ nh chunks 5 trang/call, cháº¡y song song ThreadPoolExecutor", False),
        ("Trang cÃ³ text PDF â†’ giá»¯ text + gá»­i áº£nh; trang chá»‰ cÃ³ áº£nh scan â†’ gá»­i áº£nh cho Kimi", False),
        ("Retry 1 láº§n náº¿u cold-start API tháº¥t báº¡i", False),
    ], font_size=14, bullet=True)

# Tier 2 detail
add_rect(slide, Inches(0.6), Inches(3.2), Inches(12), Inches(1.6), LIGHT_BG)
add_text(slide, Inches(0.8), Inches(3.3), Inches(3), Inches(0.4),
         "TIER 2 â€” REDUCE (GLM-5.2)", font_size=16, bold=True, color=BRAND_ORANGE)
add_multi_text(slide, Inches(0.8), Inches(3.7), Inches(11.5), Inches(1.2),
    [
        ("TrÃ­ch xuáº¥t 3 danh sÃ¡ch tá»« má»—i chunk há»£p Ä‘á»“ng:", False),
        ("  (A) Äiá»u khoáº£n loáº¡i trá»« â€” (B) KhÃ¡i niá»‡m/Ä‘á»‹nh nghÄ©a/danh má»¥c â€” (C) Háº¡n má»©c chi tráº£", False),
        ("Cháº¡y song song tá»«ng chunk, khÃ´ng bá» sÃ³t trang nÃ o", False),
    ], font_size=14, bullet=True)

# Tier 3 detail
add_rect(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1.9), LIGHT_BG)
add_text(slide, Inches(0.8), Inches(5.1), Inches(3), Inches(0.4),
         "TIER 3 â€” MERGE (GLM-5.2)", font_size=16, bold=True, color=BRAND_BLUE)
add_multi_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
    [
        ("Map tá»«ng má»¥c hÃ³a Ä‘Æ¡n vÃ o 3 danh sÃ¡ch: 2(A) trÃ¹ng tÃªn â†’ 2(B) kháº¥u trá»« giÃ¡n tiáº¿p â†’ 2(C) vÆ°á»£t háº¡n má»©c", False),
        ("LUáº¬T Cá»¨NG: khÃ´ng tá»± phÃ¢n loáº¡i, chá»‰ dá»±a vÃ o há»£p Ä‘á»“ng; má»—i má»¥c Ä‘á»™c láº­p 100%, khÃ´ng kÃ©o theo", False),
        ("LÃ­ do kháº¥u trá»« = trÃ­ch dáº«n Ä‘iá»u khoáº£n, khÃ´ng bá»• sung lÃ½ do ngoÃ i", False),
        ("Xuáº¥t 1 báº£ng duy nháº¥t: Tá»•ng tiá»n â†’ Kháº¥u trá»« â†’ Tiá»n cÃ²n láº¡i", False),
        ("Retry 1 láº§n náº¿u response rá»—ng/tháº¥t báº¡i", False),
    ], font_size=14, bullet=True)

# ============================================================
# SLIDE 6: TÃ­nh nÄƒng chÃ­nh
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "TÃ­nh nÄƒng chÃ­nh", "V.0.10")
slide_footer(slide, 6)

features = [
    ("Chat tá»± nhiÃªn", "Bot chÃ o há»i, há»i linh hoáº¡t 1-2 cÃ¢u/lÆ°á»£t nhÆ° Ä‘Ã¡nh giÃ¡ viÃªn tháº­t (AIML + AI)"),
    ("11 sáº£n pháº©m ", "Combo 360, Family Care, TNDS, chÃ¡y ná»•, ung thÆ°, tai náº¡n, hiá»ƒm nghÃ¨o..."),
    ("Äá»c áº£nh hÃ³a Ä‘Æ¡n", "Kimi K2.6 OCR â€” Ä‘á»c chÃ­nh xÃ¡c tá»«ng dÃ²ng, tá»«ng con sá»‘"),
    ("Äá»c há»£p Ä‘á»“ng PDF", "PyMuPDF + Kimi â€” xá»­ lÃ½ text + áº£nh scan, chia chunk song song"),
    ("Suy luáº­n kháº¥u trá»«", "Kháº¥u trá»« trá»±c tiáº¿p + giÃ¡n tiáº¿p + vÆ°á»£t háº¡n má»©c, cÃ³ trÃ­ch dáº«n Ä‘iá»u khoáº£n"),
    ("Chá»‘ng kÃ©o theo", "Má»—i má»¥c Ä‘á»™c láº­p 100%, khÃ´ng kháº¥u trá»« má»¥c khÃ´ng Ä‘Æ°á»£c há»£p Ä‘á»“ng nháº¯c Ä‘áº¿n"),
    ("Chá»‘ng tá»± phÃ¢n loáº¡i", "KhÃ´ng tá»± gÃ¡n nhÃ£n thuá»‘c/thá»±c pháº©m chá»©c nÄƒng, chá»‰ dá»±a vÃ o há»£p Ä‘á»“ng"),
    ("Retry cold-start", "Tá»± retry 1 láº§n náº¿u API cold-start tráº£ rá»—ng/tháº¥t báº¡i"),
    ("LÆ°u há»“ sÆ¡ tá»± Ä‘á»™ng", "JSON (claim_logs/) + Markdown, dá»… tra cá»©u"),
    ("Tra cá»©u vÄƒn phÃ²ng", "Danh sÃ¡ch  toÃ n quá»‘c, tÃ¬m theo thÃ nh phá»‘"),
]

for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 1.15)
    # Checkmark
    add_text(slide, x, y, Inches(0.4), Inches(0.4), "âœ…", font_size=18, color=GREEN)
    add_text(slide, x + Inches(0.5), y, Inches(5.3), Inches(0.4), title,
             font_size=16, bold=True, color=BRAND_BLUE)
    add_text(slide, x + Inches(0.5), y + Inches(0.4), Inches(5.3), Inches(0.6), desc,
             font_size=13, color=DARK_TEXT)

# ============================================================
# SLIDE 7: VÃ­ dá»¥ thá»±c táº¿
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "VÃ­ dá»¥ thá»±c táº¿", "Äá»‘i chiáº¿u hÃ³a Ä‘Æ¡n vá»›i há»£p Ä‘á»“ng")
slide_footer(slide, 7)

# Table
add_text(slide, Inches(0.8), Inches(1.4), Inches(12), Inches(0.4),
         "Há»£p Ä‘á»“ng ghi loáº¡i trá»« Má»¥c 10: Sanlein â†’ Kháº¥u trá»« Sanlein, KHÃ”NG kÃ©o theo Liposic",
         font_size=15, bold=True, color=BRAND_BLUE)

# Simple table representation
rows = [
    ["#", "Má»¥c bá»‹ kháº¥u trá»«", "Sá»‘ tiá»n (VNÄ)", "LÃ­ do", "Tiá»n cÃ²n láº¡i"],
    ["0", "Tá»•ng tiá»n ban Ä‘áº§u: 817.520", "-", "-", "817.520"],
    ["1", "Sanlein 0.3% 5ml (SL: 2)", "264.600", "Há»£p Ä‘á»“ng ghi rÃµ loáº¡i trá»« Sanlein (Má»¥c 10, Trang 4/7/8)", "552.920"],
    ["KQ", "Tá»”NG KHáº¤U TRá»ª", "264.600", "", "552.920"],
]

tbl_left = Inches(0.8)
tbl_top = Inches(2.0)
tbl_w = Inches(12)
tbl_h = Inches(2.5)
table_shape = slide.shapes.add_table(len(rows), 5, tbl_left, tbl_top, tbl_w, tbl_h)
table = table_shape.table

# Column widths
table.columns[0].width = Inches(0.6)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(5)
table.columns[4].width = Inches(1.4)

for ri, row_data in enumerate(rows):
    for ci, val in enumerate(row_data):
        cell = table.cell(ri, ci)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(13)
        if ri == 0:
            para.font.bold = True
            para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND_BLUE
        elif ri == len(rows) - 1:
            para.font.bold = True
            para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND_ORANGE
        else:
            para.font.color.rgb = DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 1 else WHITE

# Note about Liposic
add_rect(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(1.5), LIGHT_BG)
add_text(slide, Inches(1), Inches(5.1), Inches(11.5), Inches(0.4),
         "Liposic Eye gel 2% 10g â€” KHÃ”NG bá»‹ kháº¥u trá»«", font_size=15, bold=True, color=GREEN)
add_text(slide, Inches(1), Inches(5.5), Inches(11.5), Inches(1),
         "Há»£p Ä‘á»“ng khÃ´ng nháº¯c Ä‘áº¿n Liposic â†’ KHÃ”NG kháº¥u trá»«. Bot khÃ´ng tá»± suy 'tÆ°Æ¡ng tá»± Sanlein' hay 'cÃ¹ng lÃ  thá»±c pháº©m chá»©c nÄƒng'. "
         "Má»—i má»¥c Ä‘á»™c láº­p 100%, chá»‰ kháº¥u trá»« khi há»£p Ä‘á»“ng thá»±c sá»± ghi rÃµ.",
         font_size=13, color=DARK_TEXT)

# ============================================================
# SLIDE 8: Lá»™ trÃ¬nh phÃ¡t triá»ƒn
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, "Lá»™ trÃ¬nh phÃ¡t triá»ƒn", "CÃ¡c phiÃªn báº£n Ä‘Ã£ ship")
slide_footer(slide, 8)

versions = [
    ("V.0.1 â€” V.0.3", "Prototype", [
        "Chatbot cÆ¡ báº£n Streamlit + OpenAI",
        "Nháº­p liá»‡u 11 sáº£n pháº©m ",
        "Luá»“ng há»i Ä‘Ã¡p Ä‘Ã¡nh giÃ¡ tiáº¿p nháº­n",
    ]),
    ("V.0.4 â€” V.0.6", "Chatbot + AIML", [
        "Bot chÃ o há»i tá»± nhiÃªn vá»›i AIML",
        "LÆ°u há»“ sÆ¡ JSON tá»± Ä‘á»™ng",
        "Tra cá»©u vÄƒn phÃ²ng  toÃ n quá»‘c",
    ]),
    ("V.0.7 â€” V.0.9", "AI Deduction Pipeline", [
        "Pipeline 3 táº§ng Map-Reduce-Merge",
        "Kimi K2.6 Ä‘á»c áº£nh + GLM-5.2 phÃ¢n tÃ­ch",
        "Xá»­ lÃ½ PDF text + áº£nh scan song song",
    ]),
    ("V.0.10", "Guardrail + Retry (hiá»‡n táº¡i)", [
        "Cáº¥m tá»± phÃ¢n loáº¡i thuá»‘c/thá»±c pháº©m chá»©c nÄƒng",
        "Luáº­t cá»©ng chá»‘ng kÃ©o theo má»¥c khÃ¡c",
        "Retry cold-start API cho Tier 1 + Tier 3",
    ]),
]

for i, (ver, title, items) in enumerate(versions):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 2.7)
    add_rect(slide, x, y, Inches(5.8), Inches(2.4), LIGHT_BG)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.4),
             ver, font_size=13, bold=True, color=BRAND_ORANGE)
    add_text(slide, x + Inches(0.2), y + Inches(0.4), Inches(5.4), Inches(0.4),
             title, font_size=17, bold=True, color=BRAND_BLUE)
    add_multi_text(slide, x + Inches(0.2), y + Inches(0.9), Inches(5.4), Inches(1.5),
        [(item, False) for item in items], font_size=13, bullet=True)

# ============================================================
# SLIDE 9: Káº¿t luáº­n
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRAND_BLUE)
add_rect(slide, 0, Inches(1.5), W, Inches(4.5), RGBColor(0x00, 0x1F, 0x50))

add_text(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.8),
         "Káº¿t luáº­n", font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_multi_text(slide, Inches(2), Inches(2.8), Inches(9.3), Inches(3),
    [
        ("Bot tá»± Ä‘á»™ng Ä‘á»c hÃ³a Ä‘Æ¡n + há»£p Ä‘á»“ng, Ä‘á»‘i chiáº¿u kháº¥u trá»« chÃ­nh xÃ¡c", True),
        ("KhÃ´ng cáº§n nhÃ¢n viÃªn Ä‘á»c láº¡i há»£p Ä‘á»“ng thá»§ cÃ´ng", False),
        ("Má»—i má»¥c Ä‘á»™c láº­p, cÃ³ trÃ­ch dáº«n Ä‘iá»u khoáº£n, khÃ´ng kÃ©o theo", False),
        ("Hoáº¡t Ä‘á»™ng vá»›i má»i há»£p Ä‘á»“ng â€” khÃ´ng cáº§n tra cá»©u thuá»‘c cá»‘ Ä‘á»‹nh", False),
        ("Há»“ sÆ¡ lÆ°u tá»± Ä‘á»™ng, dá»… tra cá»©u, dá»… audit", False),
        ("", False),
        ("Sáºµn sÃ ng má»Ÿ rá»™ng: thÃªm sáº£n pháº©m, thÃªm model, thÃªm ngÃ´n ngá»¯", True),
    ], font_size=18, color=WHITE, bullet=True)

add_text(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
         "Cáº£m Æ¡n má»i ngÆ°á»i Ä‘Ã£ láº¯ng nghe ðŸ…", font_size=22, bold=True,
         color=BRAND_ORANGE, align=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = r"D:\WORK\back up 1\_AI_Claim_Chatbot_V0.10.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")

